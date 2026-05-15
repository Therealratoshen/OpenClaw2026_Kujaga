#!/usr/bin/env python3
"""
Kujaga Hackathon Presentation Generator
Creates PPTX with all slides and placeholder images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# Output directory
OUTPUT_DIR = "/Users/filberthenrico/RISTEK x Build Club OpenClaw Agenthon 2026/OpenClaw2026_Kujaga"

def create_placeholder_image(filename, width, height, text, bg_color=(15, 15, 25), text_color=(255, 255, 255)):
    """Create a placeholder image with text"""
    img = Image.new('RGB', (width, height), bg_color)
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)

    # Add text centered
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 40)
        small_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
    except:
        font = ImageFont.load_default()
        small_font = font

    # Calculate text position
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    x = (width - text_width) // 2
    y = (height - text_height) // 2
    draw.text((x, y), text, fill=text_color, font=font)

    # Add small text at bottom
    subtext = "[Placeholder Image]"
    sub_bbox = draw.textbbox((0, 0), subtext, font=small_font)
    sub_width = sub_bbox[2] - sub_bbox[0]
    draw.text(((width - sub_width) // 2, height - 50), subtext, fill=(100, 100, 100), font=small_font)

    filepath = os.path.join(OUTPUT_DIR, filename)
    img.save(filepath)
    print(f"Created: {filename}")
    return filepath

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_top = Inches(4.2)
    sub_height = Inches(1)
    sub_box = slide.shapes.add_textbox(left, sub_top, width, sub_height)
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(9)
    height = Inches(1.5)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Bullets
    bullet_left = Inches(0.5)
    bullet_top = Inches(1.3)
    bullet_width = Inches(5.5) if image_path else Inches(9)
    bullet_height = Inches(5)

    bullet_box = slide.shapes.add_textbox(bullet_left, bullet_top, bullet_width, bullet_height)
    tf2 = bullet_box.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(12)

    # Image placeholder
    if image_path:
        img_left = Inches(6)
        img_top = Inches(1.5)
        slide.shapes.add_picture(image_path, img_left, img_top, width=Inches(3.5))

def add_image_slide(prs, title, image_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.5), width=Inches(7))

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(0.5))
        tf2 = cap_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(148, 163, 184)
        p2.alignment = PP_ALIGN.CENTER

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Create placeholder images
    hero_img = create_placeholder_image("slide_hero.png", 800, 400, "KUJAGA\nAutonomous Identity Security Agent", (88, 28, 214))
    telegram_img = create_placeholder_image("slide_telegram.png", 400, 300, "Telegram Bot\n@Kujaga_bot", (0, 136, 204))
    dashboard_img = create_placeholder_image("slide_dashboard.png", 400, 300, "Dashboard Preview", (15, 15, 25))
    architecture_img = create_placeholder_image("slide_arch.png", 400, 300, "Architecture Flow", (34, 197, 94))
    demo_img = create_placeholder_image("slide_demo.png", 600, 350, "DEMO Screenshot", (15, 15, 25))

    # ============ SLIDE 1: Title ============
    add_title_slide(prs,
        "Kujaga 🛡️",
        "Autonomous Identity Security Agent\nfor Indonesian Public Figures\n\nRISTEK x Build Club OpenClaw Agenthon 2026")

    # ============ SLIDE 2: Inspiration ============
    add_content_slide(prs, "💡 Inspiration", [
        "Indonesia lost Rp 9.1 Trillion to cybercrime last year",
        "432,000+ digital fraud cases — increasing 40% YoY",
        "148 Million Indonesian citizens affected by data breaches",
        "Public figures (influencers, creators, entrepreneurs) are prime targets:",
        "   - Account cloning for financial fraud",
        "   - Impersonation for investment scams",
        "   - Data doxxing and physical threats",
        "Enterprise security costs Rp 500Juta-Rp 5Million/month",
        "Only ~35% of Indonesian creators have any digital protection",
        "\"Kujaga\" = pelindung (protector) in Sanskrit"
    ], hero_img)

    # ============ SLIDE 3: What It Does ============
    add_content_slide(prs, "🎯 What It Does", [
        "Kujaga is an Autonomous Identity Security Agent that works 24/7:",
        "",
        "🔐 Breach Detection — Monitor emails across 15B+ breach records",
        "👁️ Brand Monitoring — Track name mentions across the internet",
        "🎣 Phishing Block — Detect impersonating domains",
        "📄 PDUPA Letter Gen — AI-generated Indonesian PDP Law letters",
        "📱 Telegram Alerts — Real-time notifications in Indonesian",
        "⏰ 30-Day Tracker — Auto-escalate to Kominfo if no response",
        "",
        "Setup once. Kujaga protects forever."
    ])

    # ============ SLIDE 4: How We Built It ============
    add_content_slide(prs, "🔧 How We Built It", [
        "OpenClaw — Self-hosted AI gateway on VPS",
        "Telegram Bot (@Kujaga_bot) — Primary user channel",
        "MiniMax API — Letter generation with Indonesian context",
        "Supabase — Database + pg_cron for scheduled monitoring",
        "Next.js 14 — Dashboard UI",
        "",
        "FREE Security APIs:",
        "• HIBP (Have I Been Pwned) — Breach check",
        "• Google Alerts RSS — Brand mentions",
        "• Google News RSS — Company breach news",
        "• VirusTotal — URL safety check"
    ], architecture_img)

    # ============ SLIDE 5: Architecture ============
    add_image_slide(prs, "🏗️ Architecture", "slide_arch.png",
        "User → Telegram → OpenClaw Gateway → Kujaga Agent → Supabase + Security APIs")

    # ============ SLIDE 6: Demo ============
    add_image_slide(prs, "📱 Demo Preview", "slide_demo.png",
        "Telegram conversation with @Kujaga_bot showing breach check, brand monitoring, and PDUPA letter generation")

    # ============ SLIDE 7: Challenges ============
    add_content_slide(prs, "⚠️ Challenges We Ran Into", [
        "OpenClaw Gateway Configuration",
        "   - Default loopback-only binding required discovery",
        "   - Correct syntax: gateway.bind lan",
        "",
        "MiniMax API Integration",
        "   - Model configured but API key auth needed careful setup",
        "",
        "Stateless vs Stateful",
        "   - Balancing autonomous monitoring with stateless API design",
        "",
        "Indonesian PDP Law Context",
        "   - Generic LLMs lack Indonesian legal terminology",
        "   - Letter generation required proper legal language"
    ])

    # ============ SLIDE 8: Accomplishments ============
    add_content_slide(prs, "🏆 Accomplishments", [
        "✅ Functional autonomous agent — 24/7 monitoring, zero triggers",
        "✅ Telegram-native experience — Setup in 2 minutes",
        "✅ Free security tools — HIBP k-anonymity, Google RSS, VirusTotal",
        "✅ Complete user flow — Signup to PDUPA letter generation",
        "✅ OpenClaw Hackathon submission — Built on OpenClaw framework",
        "",
        "Tagline: \"Kujaga menjaga namamu, data kamu, dan reputasi kamu.\""
    ])

    # ============ SLIDE 9: What We Learned ============
    add_content_slide(prs, "📚 What We Learned", [
        "Single agent architecture is correct for monitoring tasks",
        "   (sequential, not parallel — based on Google Research paper)",
        "",
        "OpenClaw's personal trust model fits Kujaga perfectly",
        "",
        "Indonesian language AI needs careful prompt engineering",
        "   for proper PDP Law terminology",
        "",
        "Telegram-first is the right call",
        "   50-70M Indonesian users already on the platform"
    ])

    # ============ SLIDE 10: What's Next ============
    add_content_slide(prs, "🚀 What's Next for Kujaga", [
        "Short-term:",
        "• Full Supabase integration (users, threats, documents)",
        "• Real DOKU payment integration",
        "• API keys for all security services",
        "",
        "Long-term:",
        "• Mobile app with deeper Telegram integration",
        "• White-label for agencies/managers",
        "• Multi-language support",
        "• Partnership with Indonesian cybersecurity firms"
    ])

    # ============ SLIDE 11: Team ============
    add_content_slide(prs, "👥 Team Kujaga", [
        "Built with ❤️ for the RISTEK x Build Club OpenClaw Agenthon 2026",
        "",
        "Special thanks to:",
        "• OpenClaw Team — For the amazing agent framework",
        "• RISTEK x Build Club — For organizing this hackathon",
        "• Indonesian cybersecurity community — For inspiration",
        "",
        "GitHub: github.com/Therealratoshen/OpenClaw2026_Kujaga",
        "Telegram: @Kujaga_bot"
    ])

    # ============ SLIDE 12: Thank You ============
    add_title_slide(prs,
        "Terima Kasih! 🙏",
        "\"Kujaga menjaga namamu, data kamu, dan reputasi kamu.\"\n\n@Kujaga_bot")

    # Save
    output_path = os.path.join(OUTPUT_DIR, "Kujaga_Hackathon_Presentation.pptx")
    prs.save(output_path)
    print(f"\n✅ Presentation saved: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()